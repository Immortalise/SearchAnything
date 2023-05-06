# pip install python-pptx

from pathlib import Path
from collections import defaultdict  
from pptx import Presentation  
from typing import List, Dict, Tuple  
from nltk.tokenize import sent_tokenize  

from .parser_base import BaseParser
from utils import encode_text


class PPTXParser(BaseParser):  
    """  
    Parser for PowerPoint files (.pptx)  
    """  
    type = 'pptx'  
      
    def __init__(self, file_path: str=None, model=None, out_path: str=None) -> None:  
        super().__init__(file_path, model, out_path)  
  
    def parse(self) -> List[Dict]:  
        page_sents = self._to_sentences()  
        if not page_sents:  
            return None  
  
        self.parse_output = []  
        for pageno, sents in page_sents:  
            for sent in sents:  
                file_dict = {}  
                file_dict['title'] = None  
                file_dict['author'] = None
                file_dict['page'] = pageno  
                file_dict['content'] = sent  
                file_dict['embedding'] = encode_text(self.model, sent)  
                file_dict['file_path'] = self.file_path  
                file_dict['subject'] = None
  
                self.parse_output.append(file_dict)  
  
        return self.parse_output  
  
    def _to_sentences(self) -> List[Tuple[int, str]]:  
        """  
        Parse PowerPoint file to text [(pageno, sentence)]  
        """  
        if not self._check_format():  
            self.parse_output = None  
            return []  
  
        pptx_doc = Presentation(self.file_path)  
  
        raw_text = ""  
        page_text = []  
        pageno = 1  
        for slide in pptx_doc.slides:  
            raw_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])  
            raw_text = raw_text.replace("\n", " ")  
            page_text.append((pageno, raw_text))  
            pageno += 1  
  
        page_sents = map(lambda x: (x[0], sent_tokenize(x[1])), page_text)  
        return page_sents  
  
    @property  
    def metadata(self) -> defaultdict:  
        if not self._metadata:  
            # PPTX files don't have built-in metadata like PDFs, so we set it to empty strings  
            metadata = defaultdict(str)  
            self._metadata = metadata  
  
        return self._metadata  
  
    def _check_format(self) -> bool:  
        f_path = Path(self.file_path)  
        return f_path.exists() and f_path.suffix == '.pptx'  
